"""Generate Tax Filing Assistant Presentation"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(15, 23, 42)
BG_CARD = RGBColor(30, 41, 59)
ACCENT = RGBColor(99, 102, 241)
ACCENT2 = RGBColor(139, 92, 246)
GREEN = RGBColor(52, 211, 153)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(148, 163, 184)
ORANGE = RGBColor(251, 146, 60)
CYAN = RGBColor(34, 211, 238)
PINK = RGBColor(244, 114, 182)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, fill_color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_para(tf, text, size=16, color=WHITE, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

# ========== SLIDE 1: Title ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
# Accent bar
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_shape(s, 4, 2.2, 5.3, 0.06, ACCENT)
add_text(s, 1.5, 2.5, 10.3, 1.2, "AI-Powered Intelligent", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(s, 1.5, 3.4, 10.3, 1.2, "Tax Filing Assistant", 52, ACCENT, True, PP_ALIGN.CENTER)
add_text(s, 2, 4.5, 9.3, 0.8, "A Production-Level Full-Stack System with 7+ AI/ML Components", 22, GRAY, False, PP_ALIGN.CENTER)
add_text(s, 2, 5.5, 9.3, 0.6, "Python  •  FastAPI  •  scikit-learn  •  NLTK  •  Plotly  •  Deep Learning", 16, CYAN, False, PP_ALIGN.CENTER)
add_text(s, 2, 6.4, 9.3, 0.5, "Semester 4 — AI Coursework Presentation", 18, GRAY, False, PP_ALIGN.CENTER)

# ========== SLIDE 2: Problem Statement ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, 0.8, 0.3, 6, 0.7, "THE PROBLEM", 32, ACCENT, True)

problems = [
    ("😵", "Complex Tax System", "Indian tax law has 2 regimes, 15+ deduction sections, and changes yearly"),
    ("📄", "Manual Filing Errors", "Over 30% of self-filed returns have errors leading to notices"),
    ("💰", "Missed Deductions", "Taxpayers miss ₹15,000-₹50,000 in savings due to lack of awareness"),
    ("⏰", "Time-Consuming", "Average taxpayer spends 8-12 hours understanding and filing taxes"),
]

for i, (icon, title, desc) in enumerate(problems):
    row = i // 2
    col = i % 2
    x = 0.8 + col * 6.2
    y = 1.3 + row * 2.8
    card = add_shape(s, x, y, 5.8, 2.4, BG_CARD)
    add_text(s, x+0.3, y+0.2, 1, 0.7, icon, 36, WHITE, False)
    add_text(s, x+0.3, y+0.9, 5.2, 0.5, title, 22, WHITE, True)
    add_text(s, x+0.3, y+1.4, 5.2, 0.8, desc, 15, GRAY)

# ========== SLIDE 3: Solution Overview ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, 0.8, 0.3, 8, 0.7, "OUR SOLUTION — 7+ AI COMPONENTS", 32, ACCENT, True)

components = [
    ("A* Search", "State-Space", ACCENT),
    ("Random Forest", "Classification", GREEN),
    ("Grad. Boost", "Regression", ORANGE),
    ("Bayesian Net", "Probabilistic", CYAN),
    ("Deep MLP", "Neural Network", PINK),
    ("NLP Engine", "Chatbot", ACCENT2),
    ("GenAI", "Generative", GREEN),
    ("Heuristic", "Optimization", ORANGE),
]

for i, (name, cat, clr) in enumerate(components):
    col = i % 4
    row = i // 4
    x = 0.8 + col * 3.1
    y = 1.3 + row * 2.8
    card = add_shape(s, x, y, 2.8, 2.4, BG_CARD)
    # Color accent bar on card
    accent_bar = add_shape(s, x, y, 2.8, 0.06, clr)
    add_text(s, x+0.2, y+0.3, 2.4, 0.4, cat, 12, clr, True)
    add_text(s, x+0.2, y+0.7, 2.4, 0.6, name, 20, WHITE, True)

# ========== SLIDE 4: System Architecture ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, 0.8, 0.3, 8, 0.7, "SYSTEM ARCHITECTURE", 32, ACCENT, True)

layers = [
    ("Frontend Layer", "HTML/CSS/JS + Plotly.js — Filing Form, Chatbot, Dashboard, Simulator, Docs", CYAN, 1.2),
    ("API Layer (FastAPI)", "tax_routes  |  chat_routes  |  analysis_routes  |  simulator_routes  |  document_routes", ACCENT, 2.4),
    ("AI/ML Model Layer", "A* Search  |  RF Classifier  |  GB Regressor  |  Bayesian Net  |  DNN (MLP)  |  NLP  |  GenAI", GREEN, 3.6),
    ("Services Layer", "Tax Calculator  |  Generative Engine  |  Prompt Templates  |  Document Processor", ORANGE, 4.8),
    ("Utilities & Data", "Privacy Logger (DPDPA)  |  Explainability  |  Datasets (CSV/JSON)  |  MLOps Pipeline", PINK, 6.0),
]

for title, desc, clr, y in layers:
    card = add_shape(s, 0.8, y, 11.7, 1.0, BG_CARD)
    bar = add_shape(s, 0.8, y, 0.08, 1.0, clr)
    add_text(s, 1.2, y+0.05, 3.5, 0.45, title, 17, clr, True)
    add_text(s, 1.2, y+0.5, 10.8, 0.45, desc, 13, GRAY)

# ========== SLIDE 5: A* State-Space Search ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, 0.8, 0.3, 10, 0.7, "AI MODEL 1 — A* STATE-SPACE SEARCH", 32, ACCENT, True)

add_shape(s, 0.8, 1.2, 5.8, 5.8, BG_CARD)
tb = add_text(s, 1.1, 1.4, 5.4, 0.5, "How It Works", 22, CYAN, True)
points = [
    "• Models tax filing as a graph search problem",
    "• States = filing steps (15 steps from Start → Complete)",
    "• Edges = valid transitions with associated costs",
    "• Heuristic: h(n) = steps×0.7 + missing_data×3.0 + errors×0.5",
    "• Admissible & consistent — guarantees optimal path",
    "• Uses heapq priority queue for efficient expansion",
    "• f(n) = g(n) + h(n) — classic A* formulation",
]
txBox = add_text(s, 1.1, 2.0, 5.4, 4.5, "", 14, GRAY)
tf = txBox.text_frame
for p_text in points:
    add_para(tf, p_text, 14, GRAY, False, 8)

add_shape(s, 7.0, 1.2, 5.5, 5.8, BG_CARD)
add_text(s, 7.3, 1.4, 5.0, 0.5, "Filing Path (A* Output)", 22, GREEN, True)
steps = ["START", "Personal Info", "Income (Salary)", "HRA Calc", "Deductions 80C", "Deductions 80D", "Other Deductions", "Regime Selection", "Tax Computation", "TDS Verify", "Tax Payment", "Verification", "FILING COMPLETE ✓"]
txBox2 = add_text(s, 7.3, 2.0, 5.0, 4.5, "", 13, WHITE)
tf2 = txBox2.text_frame
for i, step in enumerate(steps):
    clr = GREEN if i == len(steps)-1 else WHITE
    add_para(tf2, f"  {i+1}. {step}" if i < len(steps)-1 else f"  ✅ {step}", 13, clr, i==len(steps)-1, 5)

# ========== SLIDE 6: ML Classification & Regression ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, 0.8, 0.3, 12, 0.7, "AI MODEL 2 & 3 — ML CLASSIFICATION & REGRESSION", 30, ACCENT, True)

# Classification card
add_shape(s, 0.8, 1.2, 5.8, 5.8, BG_CARD)
add_text(s, 1.1, 1.4, 5.4, 0.5, "Tax Regime Classifier", 22, GREEN, True)
cls_points = [
    "Algorithm: Random Forest + Gradient Boosting Ensemble",
    "Task: Predict Old vs New regime (binary classification)",
    "Features: income, deductions, HRA, investments, age, city",
    "Feature Engineering: deduction ratios, income brackets",
    "Ensemble: Average probabilities from both models",
    "Output: Regime + confidence score + feature importance",
    "Fallback: Rule-based when models unavailable",
]
txB = add_text(s, 1.1, 2.1, 5.4, 4.5, "", 13, GRAY)
tf = txB.text_frame
for p in cls_points:
    add_para(tf, f"• {p}", 13, GRAY, False, 7)

# Regression card
add_shape(s, 7.0, 1.2, 5.5, 5.8, BG_CARD)
add_text(s, 7.3, 1.4, 5.0, 0.5, "Tax Liability Predictor", 22, ORANGE, True)
reg_points = [
    "Algorithm: Gradient Boosting Regressor",
    "Task: Predict tax amount (continuous value)",
    "Feature Engineering: log(income), deduction ratios",
    "Confidence: Mean ± std from estimator ensemble",
    "Metrics: R² score, MAE, prediction intervals",
    "Training: 130+ synthetic taxpayer records",
    "Fallback: Rule-based slab calculation",
]
txB2 = add_text(s, 7.3, 2.1, 5.0, 4.5, "", 13, GRAY)
tf2 = txB2.text_frame
for p in reg_points:
    add_para(tf2, f"• {p}", 13, GRAY, False, 7)

# ========== SLIDE 7: Bayesian Network ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, 0.8, 0.3, 12, 0.7, "AI MODEL 4 — BAYESIAN NETWORK", 32, ACCENT, True)

add_shape(s, 0.8, 1.2, 5.8, 5.8, BG_CARD)
add_text(s, 1.1, 1.4, 5.4, 0.5, "Probabilistic Reasoning", 22, CYAN, True)
bay_points = [
    "Bayes' Theorem: P(A|B) = P(B|A)·P(A) / P(B)",
    "Nodes: Income Level, Deduction Level, Regime, Audit Risk",
    "24-entry Conditional Probability Table (CPT)",
    "Likelihood modifiers for risk factors:",
    "  — High deduction-to-income ratio → ×1.8",
    "  — Maxed 80C with low income → ×1.4",
    "  — Large donations (>10% income) → ×1.5",
    "  — Income > ₹50L → ×1.6",
    "Posterior = min(1.0, base_prob × likelihood_modifier)",
]
txB = add_text(s, 1.1, 2.0, 5.4, 4.8, "", 13, GRAY)
tf = txB.text_frame
for p in bay_points:
    add_para(tf, f"• {p}" if not p.startswith("  ") else p, 13, GRAY, False, 6)

add_shape(s, 7.0, 1.2, 5.5, 5.8, BG_CARD)
add_text(s, 7.3, 1.4, 5.0, 0.5, "Use Cases", 22, GREEN, True)
uses = [
    ("🔍 Audit Risk Assessment", "Calculates P(Audit | Income, Deductions, Regime) with risk level classification (Low < 3%, Medium < 8%, High ≥ 8%)"),
    ("⚖️ Regime Recommendation", "Bayesian update with HRA & home loan signals — P(OldBetter | Data) vs P(NewBetter | Data)"),
    ("📊 Transparent Output", "Shows prior probability, conditional probability, posterior probability, and all risk factors"),
]
txB2 = add_text(s, 7.3, 2.0, 5.0, 4.8, "", 14, WHITE)
tf2 = txB2.text_frame
for title, desc in uses:
    add_para(tf2, title, 16, WHITE, True, 12)
    add_para(tf2, desc, 12, GRAY, False, 4)

# ========== SLIDE 8: Deep Learning & NLP ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, 0.8, 0.3, 12, 0.7, "AI MODEL 5 & 6 — DEEP LEARNING & NLP CHATBOT", 30, ACCENT, True)

# DL card
add_shape(s, 0.8, 1.2, 5.8, 2.6, BG_CARD)
add_text(s, 1.1, 1.3, 5.4, 0.5, "Document Classifier (DNN)", 20, PINK, True)
dl_points = ["Architecture: MLP (256→128→64) + ReLU + Dropout", "7 document categories (Form16, Investment, etc.)", "Early stopping + validation split for regularization"]
txB = add_text(s, 1.1, 1.9, 5.4, 1.8, "", 13, GRAY)
tf = txB.text_frame
for p in dl_points:
    add_para(tf, f"• {p}", 13, GRAY, False, 5)

# NLP card
add_shape(s, 0.8, 4.1, 11.7, 3.2, BG_CARD)
add_text(s, 1.1, 4.2, 11, 0.5, "NLP Chatbot Pipeline", 20, CYAN, True)
add_text(s, 1.1, 4.8, 11, 0.6, "User Message → NLTK Tokenize → TF-IDF Embed (3000 features) → Attention Layer → MLP Intent Classify (128-64-32) → Entity Extract → Context Update → Response Generate", 15, WHITE)
nlp_points = ["Custom TaxTokenizer: regex for ₹ amounts, sections (80C/80D), PAN, dates", "Attention: Scaled dot-product — softmax(Q·Kᵀ/√dₖ)·V", "25+ intent categories with data augmentation during training", "Rule-based fallback when confidence < 40%"]
txB2 = add_text(s, 1.1, 5.5, 11, 2.0, "", 13, GRAY)
tf2 = txB2.text_frame
for p in nlp_points:
    add_para(tf2, f"• {p}", 13, GRAY, False, 5)

# GenAI small card
add_shape(s, 7.0, 1.2, 5.5, 2.6, BG_CARD)
add_text(s, 7.3, 1.3, 5.0, 0.5, "Generative AI Engine", 20, ACCENT2, True)
gen_points = ["Role-Based Prompting (Tax Expert persona)", "Few-Shot Prompting (example Q&A pairs)", "Chain-of-Thought step-by-step reasoning"]
txB3 = add_text(s, 7.3, 1.9, 5.0, 1.8, "", 13, GRAY)
tf3 = txB3.text_frame
for p in gen_points:
    add_para(tf3, f"• {p}", 13, GRAY, False, 5)

# ========== SLIDE 9: API Endpoints ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, 0.8, 0.3, 8, 0.7, "REST API ENDPOINTS (FastAPI)", 32, ACCENT, True)

endpoints = [
    ("POST", "/api/tax/calculate", "Calculate tax liability", GREEN),
    ("POST", "/api/tax/recommend-regime", "ML regime recommendation", GREEN),
    ("POST", "/api/tax/predict-liability", "Predict tax amount", GREEN),
    ("POST", "/api/tax/recommend-deductions", "AI deduction tips", GREEN),
    ("POST", "/api/tax/audit-risk", "Bayesian audit assessment", GREEN),
    ("POST", "/api/tax/generate-summary", "GenAI tax summary", GREEN),
    ("POST", "/api/chat/message", "Send chatbot message", CYAN),
    ("POST", "/api/analysis/filing-path", "A* search path", ORANGE),
    ("POST", "/api/documents/classify", "DL doc classification", PINK),
    ("POST", "/api/simulator/what-if", "What-if scenarios", ACCENT2),
    ("GET", "/api/model-status", "AI model health", GRAY),
    ("GET", "/api/glossary", "Tax glossary (50+ terms)", GRAY),
]

# Header
add_shape(s, 0.8, 1.1, 11.7, 0.45, ACCENT)
add_text(s, 1.0, 1.12, 1.5, 0.4, "Method", 13, WHITE, True)
add_text(s, 2.5, 1.12, 4.5, 0.4, "Endpoint", 13, WHITE, True)
add_text(s, 7.0, 1.12, 5.0, 0.4, "Description", 13, WHITE, True)

for i, (method, endpoint, desc, clr) in enumerate(endpoints):
    y = 1.65 + i * 0.44
    if i % 2 == 0:
        add_shape(s, 0.8, y-0.02, 11.7, 0.44, BG_CARD)
    add_text(s, 1.0, y, 1.5, 0.4, method, 12, clr, True)
    add_text(s, 2.5, y, 4.5, 0.4, endpoint, 12, WHITE, False)
    add_text(s, 7.0, y, 5.0, 0.4, desc, 12, GRAY, False)

# ========== SLIDE 10: Ethics & Privacy ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, 0.8, 0.3, 8, 0.7, "RESPONSIBLE AI & ETHICS", 32, ACCENT, True)

ethics = [
    ("🔍 Transparency", "Feature importance rankings, confidence scores, methodology labels, chain-of-thought reasoning for every AI decision", CYAN),
    ("🔒 Privacy (DPDPA)", "PAN/Aadhaar/email auto-redacted in logs, session-based processing, consent management, 30-day retention policy", GREEN),
    ("⚖️ Fairness", "Balanced synthetic dataset across income levels, ages, employment types, cities — no demographic bias", ORANGE),
    ("🛡️ Safety", "Rule-based fallbacks when ML fails, input validation, no external API calls — all processing local", PINK),
    ("👤 Human Oversight", "AI assists, humans decide — users can override any recommendation, step-by-step verification", ACCENT2),
    ("📋 Accountability", "MLOps model registry, performance monitoring, version tracking, known limitations disclosed", GRAY),
]

for i, (title, desc, clr) in enumerate(ethics):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.2 + row * 2.0
    card = add_shape(s, x, y, 5.8, 1.7, BG_CARD)
    bar = add_shape(s, x, y, 0.06, 1.7, clr)
    add_text(s, x+0.3, y+0.1, 5.2, 0.4, title, 18, clr, True)
    add_text(s, x+0.3, y+0.55, 5.2, 1.0, desc, 13, GRAY)

# ========== SLIDE 11: Tech Stack ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, 0.8, 0.3, 8, 0.7, "TECHNOLOGY STACK", 32, ACCENT, True)

techs = [
    ("Backend", "Python 3.9+, FastAPI, Uvicorn, Pydantic", ACCENT),
    ("ML/AI", "scikit-learn (RF, GB, MLP), NumPy, pandas", GREEN),
    ("NLP", "NLTK, TF-IDF Vectorizer, Custom Attention", CYAN),
    ("Visualization", "Plotly.js interactive charts & dashboards", ORANGE),
    ("Frontend", "HTML5, CSS3, Vanilla JavaScript (SPA)", PINK),
    ("MLOps", "Custom pipeline, model registry, monitoring", ACCENT2),
    ("Deployment", "Docker, Render (backend), Vercel (frontend)", GRAY),
    ("Data", "Synthetic CSV (130+ records), JSON configs", WHITE),
]

for i, (cat, desc, clr) in enumerate(techs):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.2 + row * 1.45
    card = add_shape(s, x, y, 5.8, 1.2, BG_CARD)
    add_text(s, x+0.3, y+0.1, 5.2, 0.4, cat, 18, clr, True)
    add_text(s, x+0.3, y+0.55, 5.2, 0.6, desc, 14, GRAY)

# ========== SLIDE 12: Thank You ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_shape(s, 0, 0, 13.333, 0.08, ACCENT)
add_shape(s, 4, 2.2, 5.3, 0.06, ACCENT)
add_text(s, 1.5, 2.5, 10.3, 1.2, "Thank You!", 54, WHITE, True, PP_ALIGN.CENTER)
add_text(s, 1.5, 3.8, 10.3, 0.8, "AI-Powered Intelligent Tax Filing Assistant", 24, ACCENT, True, PP_ALIGN.CENTER)
add_text(s, 2, 4.8, 9.3, 0.6, "7+ AI Models  •  18 API Endpoints  •  Full-Stack Production System", 18, GRAY, False, PP_ALIGN.CENTER)
add_text(s, 2, 5.8, 9.3, 0.6, "Questions?", 28, CYAN, True, PP_ALIGN.CENTER)

# Save
out_path = os.path.join(os.path.dirname(__file__), "Tax_Filing_Assistant_Presentation.pptx")
prs.save(out_path)
print(f"✅ Presentation saved to: {out_path}")
